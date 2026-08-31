import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette: Modern Defense Cyber Theme
    BG_DARK = RGBColor(10, 22, 40)        # Deep Navy / Tactical Dark
    CARD_BG = RGBColor(18, 38, 66)        # Rich Slate Navy Container
    CARD_BORDER = RGBColor(30, 64, 106)   # Subtle outline
    ACCENT_CYAN = RGBColor(0, 229, 255)   # Tech Cyan
    ACCENT_GOLD = RGBColor(255, 183, 3)   # Defense Gold
    TEXT_WHITE = RGBColor(240, 246, 252)  # High-contrast white
    TEXT_MUTED = RGBColor(148, 163, 184)  # Secondary muted grey
    CARD_HIGHLIGHT = RGBColor(13, 148, 136) # Teal green
    ACCENT_GREEN = RGBColor(34, 197, 94)  # Verification Green

    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_header(slide, slide_num_str, title_text, category_text="AI KAVACH | CYBER REASONING SYSTEM"):
        # Header Badge / Category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = f"{category_text}  •  {slide_num_str.upper()}"
        p_c.font.size = Pt(10)
        p_c.font.bold = True
        p_c.font.color.rgb = ACCENT_CYAN

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.6))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

    def draw_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2)
        return shape

    # ==========================================
    # SLIDE 1: INTRODUCTION, IDEATION & BRIEF DESCRIPTION
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s1)
    add_header(s1, "Slide 1 of 5", "CYBERLENS-KAVACH: Autonomous Self-Healing Cyber Reasoning System")

    # Left Column: Problem Overview & National Security Motivation
    draw_card(s1, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(5.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "NATIONAL DEFENSE CHALLENGE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD

    p = tf1.add_paragraph()
    p.text = "• Zero-Window Exploitation in Armed Forces Infra:\nMilitary command, radar, and tactical networks operate under strict air-gapped constraints. Traditional manual vulnerability discovery, triage, and patch deployment cycles take 15 to 45 days, exposing strategic assets to catastrophic zero-day exploitation."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

    p = tf1.add_paragraph()
    p.text = "• The Patch Reliability Dilemma:\nAutomated patches historically risk breaking mission-critical operational logic or failing under adversary evasion because they lack formal proof-of-fix verification."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

    p = tf1.add_paragraph()
    p.text = "• Edge & Air-Gapped Barrier:\nExisting cloud-reliant AI tools violate defense telemetry policies and demand massive supercomputing clusters unviable in tactical forward bases."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

    # Right Column: Proposed Idea & Solution Overview
    draw_card(s1, Inches(6.7), Inches(1.5), Inches(5.8), Inches(5.4))
    tb2 = s1.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(5.4), Inches(5.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "PROPOSED SOLUTION: CLOSED-LOOP CRS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p = tf2.add_paragraph()
    p.text = "CyberLens-Kavach is an autonomous, lightweight Cyber Reasoning System (CRS) designed for sovereign, air-gapped defense infrastructure. It unifies hybrid discovery, local SLM intelligence, and dual-gate test verification:"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

    p = tf2.add_paragraph()
    p.text = "1. Autonomous Discovery: Fuses coverage-guided fuzzers (AFL++/Atheris) and AST-aware SAST to isolate vulnerabilities and generate PoC crash payloads."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    p = tf2.add_paragraph()
    p.text = "2. Local Reasoning Engine: Employs an air-gapped, quantized Small Reasoning Model (SLM) to synthesize surgical, AST-accurate Git diff patches."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    p = tf2.add_paragraph()
    p.text = "3. Proof-of-Fix Harness: Re-runs the exploit trigger (proves immunity) AND executes complete regression test suites (proves 100% operational continuity)."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # ==========================================
    # SLIDE 2: DETAILED METHODOLOGY & WORKFLOW
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s2)
    add_header(s2, "Slide 2 of 5", "Autonomous Tri-Stage Cyber Reasoning Methodology")

    # 3 Horizontal Workflow Columns
    col_w = Inches(3.75)
    col_h = Inches(5.4)

    # Stage 1
    draw_card(s2, Inches(0.8), Inches(1.5), col_w, col_h)
    tb_s1 = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(3.35), Inches(5.0))
    tf_s1 = tb_s1.text_frame
    tf_s1.word_wrap = True

    p = tf_s1.paragraphs[0]
    p.text = "STAGE 1: HYBRID DISCOVERY"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD

    p = tf_s1.add_paragraph()
    p.text = "• Multi-Source Ingestion:\nIngests C/C++, Rust, Python, Go codebases, microservices, and network topologies."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

    p = tf_s1.add_paragraph()
    p.text = "• Static Taint & AST Scan:\nTree-sitter & Semgrep map input sources to sensitive sinks (memory buffers, SQL, shell commands)."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    p = tf_s1.add_paragraph()
    p.text = "• Dynamic Fuzzing & PoC Gen:\nCoverage-guided fuzzers generate reproducible crash payloads with exact register & stack traces."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # Stage 2
    draw_card(s2, Inches(4.78), Inches(1.5), col_w, col_h)
    tb_s2 = s2.shapes.add_textbox(Inches(4.98), Inches(1.7), Inches(3.35), Inches(5.0))
    tf_s2 = tb_s2.text_frame
    tf_s2.word_wrap = True

    p = tf_s2.paragraphs[0]
    p.text = "STAGE 2: LLM REASONING & REPAIR"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p = tf_s2.add_paragraph()
    p.text = "• AST Slice & Context Windowing:\nIsolates only relevant code branches and crash stack frames to fit within lightweight 3B–7B SLM context limits."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

    p = tf_s2.add_paragraph()
    p.text = "• Root-Cause Cyber Reasoning:\nModel evaluates memory safety, boundary invariants, and sanitization gaps."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    p = tf_s2.add_paragraph()
    p.text = "• Unified Diff Patch Generation:\nOutputs syntactically sound Git diff patches adhering strictly to project coding standards."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # Stage 3
    draw_card(s2, Inches(8.75), Inches(1.5), col_w, col_h)
    tb_s3 = s2.shapes.add_textbox(Inches(8.95), Inches(1.7), Inches(3.35), Inches(5.0))
    tf_s3 = tb_s3.text_frame
    tf_s3.word_wrap = True

    p = tf_s3.paragraphs[0]
    p.text = "STAGE 3: PROOF-OF-FIX HARNESS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    p = tf_s3.add_paragraph()
    p.text = "• Dual-Gate Sandbox Execution:\nApplies patch in an isolated sandbox and runs automated verification."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

    p = tf_s3.add_paragraph()
    p.text = "• Gate 1 (Vulnerability Elimination):\nRe-injects the PoC crash input. Asserts exit code 0 and zero crash signals."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    p = tf_s3.add_paragraph()
    p.text = "• Gate 2 (Zero Regression):\nExecutes full functional regression suite. Proves no logic breakage."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    p = tf_s3.add_paragraph()
    p.text = "• Self-Correction Loop:\nCompiler/test errors are fed back into the SLM for up to 3 automated refinement iterations."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # ==========================================
    # SLIDE 3: TECHNOLOGY STACK & SYSTEM ARCHITECTURE
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s3)
    add_header(s3, "Slide 3 of 5", "System Architecture & High-Performance Technology Stack")

    # Left Box: System Architecture Flow Diagram
    draw_card(s3, Inches(0.8), Inches(1.5), Inches(6.0), Inches(5.4))
    tb_arch = s3.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.6), Inches(5.0))
    tf_arch = tb_arch.text_frame
    tf_arch.word_wrap = True

    p = tf_arch.paragraphs[0]
    p.text = "END-TO-END AUTONOMOUS ARCHITECTURE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p = tf_arch.add_paragraph()
    p.text = "┌──────────────────────────────────────────────┐\n│  TARGET MILITARY INFRASTRUCTURE / REPO       │\n└──────────────────────┬───────────────────────┘\n                       ▼\n┌──────────────────────────────────────────────┐\n│  HYBRID INGESTION & DISCOVERY LAYER          │\n│  • SAST (Semgrep/Tree-sitter)  • DAST (Nuclei)│\n│  • AFL++/Atheris Fuzzing Engine • Zeek/Suricata│\n└──────────────────────┬───────────────────────┘\n                       ▼ [PoC Crash + Taint Traces]\n┌──────────────────────────────────────────────┐\n│  AIR-GAPPED CYBER REASONING ENGINE           │\n│  • Ollama Local Runtime  • DeepSeek/Llama 3.2│\n│  • Root Cause Triage & AST Patch Synthesizer │\n└──────────────────────┬───────────────────────┘\n                       ▼ [Synthesized Unified Diff]\n┌──────────────────────────────────────────────┐\n│  DUAL-GATE PROOF-OF-FIX TEST HARNESS         │\n│  • Gate 1: PoC Re-execution (Vulnerability 0)│\n│  • Gate 2: Full Test Suite (Regression 0)    │\n└──────────────────────┬───────────────────────┘\n                       ▼ [Verified Fix Approved]\n┌──────────────────────────────────────────────┐\n│  AUTONOMOUS HOT-PATCH / AUDIT DEPLOYMENT     │\n└──────────────────────────────────────────────┘"
    p.font.name = "Consolas"
    p.font.size = Pt(8.5)
    p.font.color.rgb = RGBColor(165, 243, 252)
    p.space_before = Pt(6)

    # Right Box: Technology Stack Breakdown & Equipment
    draw_card(s3, Inches(7.1), Inches(1.5), Inches(5.4), Inches(5.4))
    tb_tech = s3.shapes.add_textbox(Inches(7.3), Inches(1.7), Inches(5.0), Inches(5.0))
    tf_tech = tb_tech.text_frame
    tf_tech.word_wrap = True

    p = tf_tech.paragraphs[0]
    p.text = "TECHNOLOGY STACK SPECIFICATIONS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD

    tech_items = [
        ("Reasoning AI Engine", "Ollama / llama3.2:3b / Qwen2.5-Coder (GGUF quantized for zero-cloud tactical edge deployment)"),
        ("Dynamic & Fuzzing Tools", "AFL++, Atheris, LibFuzzer, Nuclei DAST, Metasploit RPC Bridge"),
        ("Static Analysis & AST", "Semgrep OSS, Tree-sitter AST parser, Joern code property graphs"),
        ("Network & Threat Telemetry", "Suricata IDS (EVE JSON), Zeek NSM (conn/dns/http streams)"),
        ("Verification Harness", "Isolated sandbox test runners (PyTest, CTest, Cargo Test, Docker/gVisor)"),
        ("Backend & Interface", "FastAPI (Async Python), MCP Protocol, React 18 + Vite Cyber Dashboard"),
        ("Hardware / Equipment", "Runs locally on standard military laptop or 1U edge server (Min: 16GB RAM, 4 CPU cores, optional 6GB VRAM GPU)"),
    ]

    for title, desc in tech_items:
        p = tf_tech.add_paragraph()
        p.text = f"• {title}: {desc}"
        p.font.size = Pt(9.5)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(4)

    # ==========================================
    # SLIDE 4: SALIENT FEATURES, NOVELTY & USP
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s4)
    add_header(s4, "Slide 4 of 5", "Key Features, Innovation & Defense Advantages (USP)")

    # 4 Feature Cards (2x2 Grid)
    card_w = Inches(5.6)
    card_h = Inches(2.55)

    # Card 1: Proof-of-Fix
    draw_card(s4, Inches(0.8), Inches(1.5), card_w, card_h)
    tb = s4.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(2.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "1. DUAL-GATE PROOF-OF-FIX HARNESS (CORE USP)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "• Eliminates the fatal flaw of generic AI coding tools (hallucinated or broken patches).\n• Empirically proves fix holding: verifies vulnerability mitigation against live exploit payloads while strictly ensuring zero regression in existing test suites."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(4)

    # Card 2: 100% Air-Gapped & Sovereign
    draw_card(s4, Inches(6.9), Inches(1.5), card_w, card_h)
    tb = s4.shapes.add_textbox(Inches(7.1), Inches(1.6), Inches(5.2), Inches(2.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2. 100% AIR-GAPPED & ZERO CLOUD DEPENDENCY"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p = tf.add_paragraph()
    p.text = "• Tailored for Indian Armed Forces data sovereignty regulations.\n• All telemetry, models, AST analyzers, and fuzzers execute entirely on-premise without a single external byte or telemetry leak."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(4)

    # Card 3: Lightweight Resource Footprint
    draw_card(s4, Inches(0.8), Inches(4.3), card_w, card_h)
    tb = s4.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(5.2), Inches(2.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "3. HYPER-LIGHTWEIGHT RESOURCE UTILIZATION"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p = tf.add_paragraph()
    p.text = "• Optimized 3B–7B quantized Small Language Models with context pruning.\n• Operates seamlessly on edge field hardware and embedded battlefield servers, avoiding high-power GPU data centers."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(4)

    # Card 4: Multi-Layer Telemetry
    draw_card(s4, Inches(6.9), Inches(4.3), card_w, card_h)
    tb = s4.shapes.add_textbox(Inches(7.1), Inches(4.4), Inches(5.2), Inches(2.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "4. CROSS-LAYER CODE & NETWORK CORRELATION"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p = tf.add_paragraph()
    p.text = "• Integrates source AST inspection with Suricata IDS & Zeek NSM feeds.\n• Correlates active network intrusion attempts directly to vulnerable lines of source code to prioritize high-risk mission vulnerabilities."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(4)

    # ==========================================
    # SLIDE 5: FINAL DELIVERABLES & PERFORMANCE OBJECTIVES
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(s5)
    add_header(s5, "Slide 5 of 5", "Project Deliverables, Performance Objectives & PoC")

    # Left Column: Deliverables & PoC Demo
    draw_card(s5, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4))
    tb_d1 = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(5.0))
    tf_d1 = tb_d1.text_frame
    tf_d1.word_wrap = True

    p = tf_d1.paragraphs[0]
    p.text = "PROTOTYPE DELIVERABLES & POC"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    p = tf_d1.add_paragraph()
    p.text = "1. End-to-End Working CRS Application:\nModular web console and headless CLI ready for direct deployment in simulated Armed Forces testbeds."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    p = tf_d1.add_paragraph()
    p.text = "2. Autonomous Self-Repair Engine:\nReal-time code scanner with automated unified diff patch synthesis and multi-turn refinement."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    p = tf_d1.add_paragraph()
    p.text = "3. Automated Regression Harness Runner:\nDocker/gVisor sandboxed test harness outputting verifiable proof-of-fix certificates."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    p = tf_d1.add_paragraph()
    p.text = "4. Defense Audit & Executive Reporting:\nAutomated generation of defense compliance reports, CVE/CWE mapping, and attack flow graphs."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(6)

    # Right Column: Quantitative Target Metrics
    draw_card(s5, Inches(6.7), Inches(1.5), Inches(5.8), Inches(5.4))
    tb_d2 = s5.shapes.add_textbox(Inches(6.9), Inches(1.7), Inches(5.4), Inches(5.0))
    tf_d2 = tb_d2.text_frame
    tf_d2.word_wrap = True

    p = tf_d2.paragraphs[0]
    p.text = "QUANTITATIVE PERFORMANCE OBJECTIVES"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD

    metrics = [
        ("Mean Time to Patch (MTTP)", "< 90 Seconds (from crash trigger to verified patch)"),
        ("False-Positive Reduction", "> 92% reduction via dynamic PoC execution gating"),
        ("Proof-of-Fix Reliability", "100% exploit mitigation on verified patches"),
        ("Regression Test Integrity", "0% regression breakage on validated suites"),
        ("Operational Hardware Footprint", "Runs on standard 16GB RAM laptops / Edge nodes"),
        ("Network Air-Gap Compliance", "100% offline, zero external telemetry dependency"),
        ("Grand Finale Readiness", "36-hr finale-ready to run autonomously on Armed Forces custom simulated test environments"),
    ]

    for m_title, m_val in metrics:
        p = tf_d2.add_paragraph()
        p.text = f"• {m_title}:\n  {m_val}"
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(5)

    output_path = "e:/CyberLens/CYBERLENS--main/AI_Kavach_CyberLens_Submission.pptx"
    prs.save(output_path)
    print(f"Presentation successfully generated at {output_path}")

if __name__ == "__main__":
    create_deck()
